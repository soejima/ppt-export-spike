#original copy & paste extracted from: https://python-pptx.readthedocs.io/en/latest/user/charts.html
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

slide.shapes.title.text = 'Adding a Chart'

# define chart data ---------------------
chart_data = ChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (50.0, 21.4, 16.7))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

# adding a table slide
title_only_slide_layout = prs.slide_layouts[5]
slide2 = prs.slides.add_slide(title_only_slide_layout)
shapes = slide2.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('ppt-files/chart-01.pptx')
